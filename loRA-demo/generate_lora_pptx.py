#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成 LoRA 模型微调全流程 PPTX 演示文稿

用法：
    pip install python-pptx
    python generate_lora_pptx.py

输出：LoRA_微调全流程讲解.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色主题 ──
COLOR_PRIMARY     = RGBColor(0x1A, 0x56, 0xDB)   # 蓝色主色
COLOR_SECONDARY   = RGBColor(0x6C, 0x5C, 0xE7)   # 紫色辅助
COLOR_ACCENT      = RGBColor(0x10, 0xB9, 0x81)   # 绿色强调
COLOR_DARK        = RGBColor(0x1E, 0x29, 0x3B)   # 深色文字
COLOR_GRAY        = RGBColor(0x94, 0xA3, 0xB8)   # 灰色
COLOR_LIGHT_BG    = RGBColor(0xF8, 0xFA, 0xFC)   # 浅色背景
COLOR_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)   # 橙色
COLOR_RED         = RGBColor(0xEF, 0x44, 0x44)   # 红色

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ════════════════════════════════════════════════
# 工具函数
# ════════════════════════════════════════════════

def add_bg(slide, color=COLOR_WHITE):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=COLOR_DARK, spacing=Pt(6), font_name="Microsoft YaHei",
                    bullet_char="▸"):
    """添加带项目符号的列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
    return txBox


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    """添加代码块样式"""
    shape = add_rounded_rect(slide, left, top, width, height, fill_color=RGBColor(0x1E, 0x29, 0x3B))
    shape.shadow.inherit = False

    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1),
                                      width - Inches(0.4), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        p.font.name = "Consolas"
        p.space_after = Pt(1)
    return txBox


def new_slide():
    """创建新幻灯片"""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    return slide


def add_section_header(slide, section_num, title, subtitle=""):
    """添加章节标题"""
    # 左侧色条
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, fill_color=COLOR_PRIMARY)
    # 编号
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(1.5), Inches(0.8),
                 f"Part {section_num}", font_size=14, color=COLOR_GRAY, bold=True)
    # 标题
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(1.2),
                 title, font_size=40, color=COLOR_DARK, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(2.9), Inches(10), Inches(0.8),
                     subtitle, font_size=18, color=COLOR_GRAY)


def add_card(slide, left, top, width, height, title, body_lines, icon="📌",
             title_color=COLOR_PRIMARY, bg_color=COLOR_WHITE):
    """添加卡片"""
    # 卡片背景
    card = add_rounded_rect(slide, left, top, width, height, fill_color=bg_color)
    card.shadow.inherit = False

    # 顶部分隔线
    add_shape(slide, left + Inches(0.15), top + Inches(0.05),
              width - Inches(0.3), Pt(3), fill_color=title_color)

    # 标题
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.5),
                 f"{icon} {title}", font_size=15, color=title_color, bold=True)

    # 内容
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.7),
                    width - Inches(0.4), height - Inches(0.9),
                    body_lines, font_size=12, color=COLOR_DARK, spacing=Pt(4))
    return card


def add_page_number(slide, num, total):
    """添加页码"""
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=COLOR_GRAY,
                 alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════
# 幻灯片页数
# ════════════════════════════════════════════════
TOTAL_SLIDES = 24

# ────────────────────────────────────
# S1: 封面
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_DARK)

# 装饰元素
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), fill_color=COLOR_PRIMARY)
add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), fill_color=COLOR_SECONDARY)

# 主标题
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "LoRA 模型微调", font_size=54, color=COLOR_WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.8),
             "Low-Rank Adaptation — 全流程讲解", font_size=28, color=COLOR_PRIMARY, bold=True)

# 分隔线
add_shape(slide, Inches(1), Inches(3.8), Inches(3), Pt(3), fill_color=COLOR_ACCENT)

# 副信息
info_items = [
    "深入原理 · 完整流程 · 实战代码 · 测试部署",
    "从理论到生产环境的 LoRA 微调全链路指南",
    "支持 Qwen / LLaMA / ChatGLM 等主流模型架构"
]
for i, t in enumerate(info_items):
    add_text_box(slide, Inches(1), Inches(4.2 + i * 0.45), Inches(11), Inches(0.4),
                 f"✦  {t}", font_size=16, color=COLOR_GRAY)

# 底部标签
add_text_box(slide, Inches(1), Inches(6.5), Inches(5), Inches(0.4),
             "2025 | LoRA Fine-tuning Demo", font_size=12, color=COLOR_GRAY)
add_page_number(slide, 1, TOTAL_SLIDES)


# ────────────────────────────────────
# S2: 目录
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.6),
             "📋 目录", font_size=36, color=COLOR_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Pt(3), fill_color=COLOR_PRIMARY)

toc_items = [
    ("01", "LoRA 概述与核心原理", "低秩分解、数学推导、设计哲学"),
    ("02", "环境搭建与数据准备", "依赖安装、硬件要求、数据格式与预处理"),
    ("03", "模型加载与 LoRA 配置", "4bit 量化、PEFT 配置、参数选择"),
    ("04", "训练流程详解", "超参数、Trainer、损失曲线解读"),
    ("05", "模型保存与推理", "适配器保存/加载、权重合并、推理测试"),
    ("06", "评估方法与测试方案", "自动指标、人工评估、单元/集成/压力测试"),
    ("07", "生产环境部署", "vLLM/TGI/Docker/K8s、监控运维"),
    ("08", "进阶话题与思考题", "LoRA 变体、常见陷阱、深度讨论"),
]

for i, (num, title, desc) in enumerate(toc_items):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.7 + row * 1.3)

    add_shape(slide, x, y, Inches(5.8), Inches(1.1),
              fill_color=RGBColor(0xF1, 0xF5, 0xF9))

    add_text_box(slide, x + Inches(0.2), y + Inches(0.08), Inches(0.6), Inches(0.4),
                 num, font_size=24, color=COLOR_PRIMARY, bold=True)

    add_text_box(slide, x + Inches(0.8), y + Inches(0.08), Inches(4.8), Inches(0.4),
                 title, font_size=18, color=COLOR_DARK, bold=True)

    add_text_box(slide, x + Inches(0.8), y + Inches(0.55), Inches(4.8), Inches(0.4),
                 desc, font_size=12, color=COLOR_GRAY)

add_page_number(slide, 2, TOTAL_SLIDES)


# ────────────────────────────────────
# S3: Part 1 章节页
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, fill_color=COLOR_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(3), Inches(0.6),
             "PART 01", font_size=16, color=COLOR_PRIMARY, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
             "LoRA 概述与核心原理", font_size=42, color=COLOR_WHITE, bold=True)
add_shape(slide, Inches(1.5), Inches(3.9), Inches(3), Pt(3), fill_color=COLOR_ACCENT)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
             "了解 LoRA 为什么能高效微调大模型，掌握其数学原理与设计哲学",
             font_size=18, color=COLOR_GRAY)
add_page_number(slide, 3, TOTAL_SLIDES)


# ────────────────────────────────────
# S4: LoRA 核心思想
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 1.1, "LoRA 核心思想", "Low-Rank Adaptation of Large Language Models (2021)")

# 左侧：核心概念
add_card(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(3.2),
         "核心思想", [
             "冻结预训练模型全部权重",
             "在模型旁插入低秩分解矩阵 (B × A)",
             "仅更新这些少量参数完成下游任务",
             f"典型参数量节省：99%+",
         ], icon="💡")

# 右侧：数学公式
formula_box = add_rounded_rect(slide, Inches(6.8), Inches(3.8), Inches(5.8), Inches(3.2),
                                fill_color=RGBColor(0xF1, 0xF5, 0xF9))
add_text_box(slide, Inches(7.0), Inches(3.9), Inches(5.4), Inches(0.4),
             "数学形式", font_size=14, color=COLOR_PRIMARY, bold=True)
add_text_box(slide, Inches(7.0), Inches(4.4), Inches(5.4), Inches(2.2),
             "W = W₀ + ΔW\n\n"
             "ΔW = B × A\n\n"
             "B ∈ ℝᵈˣʳ, A ∈ ℝʳˣᵏ\n\n"
             "r ≪ min(d, k)       例: d=4096, r=8",
             font_size=16, color=COLOR_DARK, bold=False)
add_text_box(slide, Inches(7.0), Inches(6.5), Inches(5.4), Inches(0.4),
             "前向传播: h = W₀x + (α/r) · BAx",
             font_size=13, color=COLOR_ACCENT, bold=True)

# 顶部对比
add_text_box(slide, Inches(0.8), Inches(3.1), Inches(12), Inches(0.5),
             '▸ 为什么叫 Low-Rank？—— 预训练模型权重变化的"固有秩"很低，因此可用低维子空间近似',
             font_size=14, color=COLOR_GRAY)
add_page_number(slide, 4, TOTAL_SLIDES)


# ────────────────────────────────────
# S5: LoRA vs 全量微调
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 1.2, "LoRA vs 全量微调", "计算效率与参数效率的全面对比")

# 表格 header
headers = ["对比维度", "全量微调 (Full FT)", "LoRA 微调", "优势"]
rows = [
    ["可训练参数", "全部 (100%)", "0.01% ~ 1%", "LoRA 🏆"],
    ["显存需求 (7B)", "~56 GB", "~16 GB (4bit)", "LoRA 🏆"],
    ["训练时间", "长", "缩短 60-80%", "LoRA 🏆"],
    ["多任务部署", "N × 全量模型", "1 个基座 + N 个适配器", "LoRA 🏆"],
    ["灾难性遗忘", "高风险", "低风险（冻结基座）", "LoRA 🏆"],
    ["效果上限", "高（数据充足时）", "接近全量微调", "持平"],
    ["实现复杂度", "简单", "中等（需调 r/α）", "持平"],
]

# 表头
y_start = Inches(3.5)
col_widths = [Inches(2.0), Inches(3.0), Inches(3.5), Inches(2.5)]
col_starts = [Inches(0.8)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

for j, (h, cs, cw) in enumerate(zip(headers, col_starts, col_widths)):
    add_text_box(slide, cs + Inches(0.1), y_start, cw - Inches(0.2), Inches(0.4),
                 h, font_size=13, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, cs, y_start, cw, Inches(0.4), fill_color=COLOR_PRIMARY)

# 数据行
for i, row_data in enumerate(rows):
    y = y_start + Inches(0.4) + i * Inches(0.42)
    bg = RGBColor(0xF8, 0xFA, 0xFC) if i % 2 == 0 else COLOR_WHITE
    for j, (cell, cs, cw) in enumerate(zip(row_data, col_starts, col_widths)):
        cell_color = COLOR_PRIMARY if j == len(row_data) - 1 and "LoRA" in cell else COLOR_DARK
        cell_bold = j == len(row_data) - 1 and "LoRA" in cell
        add_shape(slide, cs, y, cw, Inches(0.42), fill_color=bg)
        add_text_box(slide, cs + Inches(0.1), y + Inches(0.02), cw - Inches(0.2), Inches(0.38),
                     cell, font_size=12, color=cell_color, bold=cell_bold,
                     alignment=PP_ALIGN.CENTER)

add_page_number(slide, 5, TOTAL_SLIDES)


# ────────────────────────────────────
# S6: 参数量计算示例
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 1.3, "参数量计算示例", "直观感受 LoRA 的效率优势")

add_card(slide, Inches(0.8), Inches(3.5), Inches(3.8), Inches(3.2),
         "假设条件", [
             "模型: d=4096, k=4096",
             "原始权重: 4096×4096 = 16.8M",
             "LoRA r=8",
             "B: 4096×8,  A: 8×4096",
         ], icon="📐")

add_card(slide, Inches(5.0), Inches(3.5), Inches(3.8), Inches(3.2),
         "计算过程", [
             "ΔW 参数量 = 4096×8 + 8×4096",
             "= 32,768 + 32,768",
             "= 65,536",
             "节省 = 1 - 65,536/16,777,216",
         ], icon="🧮", title_color=COLOR_SECONDARY)

add_card(slide, Inches(9.2), Inches(3.5), Inches(3.4), Inches(3.2),
         "结论", [
             "节省比例: 99.61%",
             "从 16.8M → 65K",
             "相当于仅用 0.39% 参数",
             "注意: 推理无额外开销！",
         ], icon="🎯", title_color=COLOR_ACCENT)

add_text_box(slide, Inches(0.8), Inches(3.0), Inches(12), Inches(0.4),
             "▸ 即使是 r=1 的极端情况，在某些任务上也能达到不错的效果",
             font_size=14, color=COLOR_GRAY)
add_page_number(slide, 6, TOTAL_SLIDES)


# ────────────────────────────────────
# S7: PART 2 章节页
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, fill_color=COLOR_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(3), Inches(0.6),
             "PART 02", font_size=16, color=COLOR_PRIMARY, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
             "环境搭建与数据准备", font_size=42, color=COLOR_WHITE, bold=True)
add_shape(slide, Inches(1.5), Inches(3.9), Inches(3), Pt(3), fill_color=COLOR_ACCENT)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
             "快速搭建微调环境，准备高质量训练数据",
             font_size=18, color=COLOR_GRAY)
add_page_number(slide, 7, TOTAL_SLIDES)


# ────────────────────────────────────
# S8: 环境依赖
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 2.1, "环境依赖与安装")

add_card(slide, Inches(0.8), Inches(3.2), Inches(5.8), Inches(3.8),
         "核心库", [
             "Python ≥ 3.10",
             "PyTorch ≥ 2.0.0",
             "transformers ≥ 4.30.0",
             "peft ≥ 0.5.0 （HuggingFace PEFT）",
             "datasets, accelerate, bitsandbytes",
             "trl (DPO/PPO 训练，可选)",
             "wandb / tensorboard（实验追踪）",
         ], icon="📦")

add_card(slide, Inches(7.0), Inches(3.2), Inches(5.6), Inches(3.8),
         "安装命令", [
             "pip install torch torchvision torchaudio",
             "pip install transformers datasets peft",
             "pip install accelerate bitsandbytes",
             "",
             "可选:",
             "pip install wandb tensorboard",
             "pip install trl vllm fastapi",
         ], icon="⚡", title_color=COLOR_SECONDARY)

# 硬件速查表
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(3), Inches(0.4),
             "硬件要求速查", font_size=13, color=COLOR_GRAY)

add_code_block(slide, Inches(0.8), Inches(7.2), Inches(11.8), Inches(0.8),
               "7B 模型 LoRA 微调: ~16GB 显存 (4bit)  |  70B 模型: ~80GB 显存 (4bit, 需 A100)",
               font_size=11)
add_page_number(slide, 8, TOTAL_SLIDES)


# ────────────────────────────────────
# S9: 数据准备
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 2.2, "数据格式与预处理流程")

# 左侧：数据格式
add_card(slide, Inches(0.8), Inches(3.2), Inches(5.8), Inches(3.0),
         "SFT 数据格式 (JSONL)", [
             '{"instruction": "...", "input": "", "output": "..."}',
             ' 或 ChatML 格式',
             '{"messages": [',
             '  {"role": "user", "content": "..."},',
             '  {"role": "assistant", "content": "..."}',
             ']}',
         ], icon="📄")

# 右侧：预处理
add_card(slide, Inches(7.0), Inches(3.2), Inches(5.6), Inches(3.0),
         "预处理流程", [
             "原始数据 → 清洗去重 → 格式化",
             "→ Tokenization → 打包/切割 → Dataset",
             "",
             "关键注意点:",
             "• loss 掩码：仅输出部分计算 loss",
             "• 质量 > 数量：500条优质 > 50000条噪声",
             "• 数据混合：通用+领域 防灾难性遗忘",
         ], icon="🔄", title_color=COLOR_SECONDARY)

add_text_box(slide, Inches(0.8), Inches(2.7), Inches(12), Inches(0.4),
             "▸ 推荐使用 JSONL 格式，每行一个完整 JSON 对象，便于分布式处理和流式读取",
             font_size=14, color=COLOR_GRAY)

# 代码示例
code = '''def preprocess_function(examples, tokenizer, max_length=2048):
    """构建 prompt，对输出部分计算 loss"""
    prompts = [f"指令: {inst}\\n输出:" for inst in examples["instruction"]]
    model_inputs = tokenizer(prompts, max_length=max_length, truncation=True)
    labels = tokenizer([p+o for p,o in zip(prompts, examples["output"])],
                       max_length=max_length, truncation=True)
    model_inputs["labels"] = labels["input_ids"]
    return model_inputs'''

add_code_block(slide, Inches(0.8), Inches(6.5), Inches(11.8), Inches(0.8), code, font_size=10)
add_page_number(slide, 9, TOTAL_SLIDES)


# ────────────────────────────────────
# S10: PART 3 章节页
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, fill_color=COLOR_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(3), Inches(0.6),
             "PART 03", font_size=16, color=COLOR_PRIMARY, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
             "模型加载与 LoRA 配置", font_size=42, color=COLOR_WHITE, bold=True)
add_shape(slide, Inches(1.5), Inches(3.9), Inches(3), Pt(3), fill_color=COLOR_ACCENT)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
             "4bit 量化加载、PEFT LoRA 配置参数详解",
             font_size=18, color=COLOR_GRAY)
add_page_number(slide, 10, TOTAL_SLIDES)


# ────────────────────────────────────
# S11: 模型加载
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 3.1, "模型加载（支持 4bit 量化）")

add_code_block(slide, Inches(0.8), Inches(3.2), Inches(11.8), Inches(2.8),
'''from transformers import AutoModelForCausalLM, BitsAndBytesConfig
import torch

# 4bit 量化配置
bnb_config = BitsAndBytesConfig(
    load_in_4bit=True,
    bnb_4bit_compute_dtype=torch.bfloat16,
    bnb_4bit_quant_type="nf4",     # NormalFloat4
    bnb_4bit_use_double_quant=True,# 双重量化
)

model = AutoModelForCausalLM.from_pretrained(
    "Qwen/Qwen2.5-7B-Instruct",
    quantization_config=bnb_config,
    device_map="auto",
    torch_dtype=torch.bfloat16,
)
model = prepare_model_for_kbit_training(model)''', font_size=11)

# 提示
add_card(slide, Inches(0.8), Inches(6.3), Inches(5.8), Inches(1.0),
         "NF4 vs FP4", [
             "NF4 假设权重正态分布，用分位数均匀分配量化级别",
             "量化误差期望最小化，比普通 4bit 精度更高",
         ], icon="💡")

add_card(slide, Inches(7.0), Inches(6.3), Inches(5.6), Inches(1.0),
         "双重量化", [
             "对量化常数再做一次量化",
             "额外节省约 0.5 bit/参数",
         ], icon="💡", title_color=COLOR_SECONDARY)
add_page_number(slide, 11, TOTAL_SLIDES)


# ────────────────────────────────────
# S12: LoRA 配置参数
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 3.2, "LoRA 配置参数详解")

add_code_block(slide, Inches(0.8), Inches(3.2), Inches(11.8), Inches(1.3),
'''from peft import LoraConfig, get_peft_model

lora_config = LoraConfig(
    r=8,                    # 秩（越小参数越少）
    lora_alpha=16,          # 缩放因子 ΔW = (α/r)·BA
    target_modules=["q_proj", "v_proj", "k_proj", "o_proj",
                    "gate_proj", "up_proj", "down_proj"],
    lora_dropout=0.05,      # 防止过拟合
    bias="none",
    task_type="CAUSAL_LM",
    use_rslora=True,        # Rank-Stabilized LoRA
    use_dora=False,         # Weight-Decomposed LoRA
)
model = get_peft_model(model, lora_config)
model.print_trainable_parameters()
# Output: trainable params: 8.3M || all params: 6.7B || trainable%: 0.12''', font_size=11)

# 参数表格
params_info = [
    ("r", "8-16", "通用 8, 复杂 16, 简单分类 4"),
    ("lora_alpha", "16-32", "通常设为 2×r"),
    ("target_modules", "全线性层", "LLM 所有线性层效果最好"),
    ("lora_dropout", "0.0-0.1", "数据少设 0.1，数据多设 0.0"),
]

y = Inches(4.8)
for i, (param, default, note) in enumerate(params_info):
    bg = RGBColor(0xF1, 0xF5, 0xF9) if i % 2 == 0 else COLOR_WHITE
    add_shape(slide, Inches(0.8), y + i * Inches(0.45), Inches(2.0), Inches(0.4), fill_color=bg)
    add_shape(slide, Inches(2.8), y + i * Inches(0.45), Inches(2.5), Inches(0.4), fill_color=bg)
    add_shape(slide, Inches(5.3), y + i * Inches(0.45), Inches(7.3), Inches(0.4), fill_color=bg)
    add_text_box(slide, Inches(0.9), y + i * Inches(0.45), Inches(1.8), Inches(0.4),
                 param, font_size=13, color=COLOR_PRIMARY, bold=True)
    add_text_box(slide, Inches(2.9), y + i * Inches(0.45), Inches(2.3), Inches(0.4),
                 default, font_size=13, color=COLOR_DARK)
    add_text_box(slide, Inches(5.4), y + i * Inches(0.45), Inches(7.1), Inches(0.4),
                 note, font_size=13, color=COLOR_GRAY)
add_page_number(slide, 12, TOTAL_SLIDES)


# ────────────────────────────────────
# S13: PART 4 章节页
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, fill_color=COLOR_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(3), Inches(0.6),
             "PART 04", font_size=16, color=COLOR_PRIMARY, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
             "训练流程详解", font_size=42, color=COLOR_WHITE, bold=True)
add_shape(slide, Inches(1.5), Inches(3.9), Inches(3), Pt(3), fill_color=COLOR_ACCENT)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
             "超参数配置、Trainer 使用、损失曲线解读与陷阱排查",
             font_size=18, color=COLOR_GRAY)
add_page_number(slide, 13, TOTAL_SLIDES)


# ────────────────────────────────────
# S14: 训练参数
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 4.1, "训练核心参数配置")

add_code_block(slide, Inches(0.8), Inches(3.2), Inches(11.8), Inches(2.5),
'''training_args = TrainingArguments(
    output_dir="./lora_output",
    num_train_epochs=3,
    per_device_train_batch_size=4,
    gradient_accumulation_steps=8,   # 有效 batch = 4×8 = 32
    learning_rate=2e-4,              # LoRA 常用 LR（比全量 FT 大 10×）
    warmup_ratio=0.03,
    lr_scheduler_type="cosine",
    bf16=True,
    optim="paged_adamw_8bit",        # 8bit 优化器省显存
    gradient_checkpointing=True,     # 省显存换速度
    weight_decay=0.01,
    max_grad_norm=1.0,
    save_strategy="steps",
    save_steps=500,
    save_total_limit=3,
)''', font_size=11)

# 关键参数说明
add_card(slide, Inches(0.8), Inches(6.0), Inches(3.8), Inches(1.3),
         "学习率策略", [
             "LoRA LR = 2e-4 (比全量FT大10-100×)",
             "原因：更新的参数少，每个参数需要走更大的步",
         ], icon="🎯")
add_card(slide, Inches(5.0), Inches(6.0), Inches(3.8), Inches(1.3),
         "梯度累积", [
             "有效 batch = per_device × grad_acc",
             "大 batch 稳定训练，但需注意显存",
         ], icon="⚙️", title_color=COLOR_SECONDARY)
add_card(slide, Inches(9.2), Inches(6.0), Inches(3.4), Inches(1.3),
         "梯度检查点", [
             "节省显存 ~40%",
             "训练速度降低 ~20%",
         ], icon="💾", title_color=COLOR_ACCENT)
add_page_number(slide, 14, TOTAL_SLIDES)


# ────────────────────────────────────
# S15: 损失曲线解读
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 4.2, "损失曲线解读指南", "通过 loss 曲线判断训练状态")

loss_patterns = [
    ("持续下降，最终收敛", "✅ 正常训练", RGBColor(0x10, 0xB9, 0x81)),
    ("快速下降到接近 0", "⚠ 过拟合 / 数据泄露", COLOR_ORANGE),
    ("震荡不收敛", "❌ 学习率过大 / 数据有冲突", COLOR_RED),
    ("loss 不降", "❌ 学习率过小 / 模型冻结", COLOR_RED),
    ("验证 loss 上升", "⚠ 过拟合，需早停 / 加正则化", COLOR_ORANGE),
]

for i, (pattern, diagnosis, color) in enumerate(loss_patterns):
    y = Inches(3.5 + i * 0.7)
    shape = add_rounded_rect(slide, Inches(0.8), y, Inches(11.8), Inches(0.55),
                              fill_color=RGBColor(0xF8, 0xFA, 0xFC))
    add_shape(slide, Inches(0.8), y, Pt(5), Inches(0.55), fill_color=color)
    add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(5), Inches(0.45),
                 pattern, font_size=14, color=COLOR_DARK, bold=True)
    add_text_box(slide, Inches(6.5), y + Inches(0.05), Inches(6), Inches(0.45),
                 diagnosis, font_size=14, color=color, bold=True)

add_text_box(slide, Inches(0.8), Inches(7.0), Inches(12), Inches(0.4),
             "💡 建议：训练初期密切关注 loss，若异常立即终止排查，避免浪费算力",
             font_size=14, color=COLOR_GRAY)
add_page_number(slide, 15, TOTAL_SLIDES)


# ────────────────────────────────────
# S16: PART 5 章节页
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, fill_color=COLOR_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(3), Inches(0.6),
             "PART 05", font_size=16, color=COLOR_PRIMARY, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
             "模型保存与推理", font_size=42, color=COLOR_WHITE, bold=True)
add_shape(slide, Inches(1.5), Inches(3.9), Inches(3), Pt(3), fill_color=COLOR_ACCENT)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
             "适配器保存与加载、权重合并、推理测试",
             font_size=18, color=COLOR_GRAY)
add_page_number(slide, 16, TOTAL_SLIDES)


# ────────────────────────────────────
# S17: 保存与合并
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 5.1, "LoRA 适配器保存 / 加载 / 合并")

# 方式一
add_card(slide, Inches(0.8), Inches(3.2), Inches(5.8), Inches(1.8),
         "方式一：保存 LoRA 适配器（推荐）", [
             "model.save_pretrained(\"./lora_adapter\")",
             "tokenizer.save_pretrained(\"./lora_adapter\")",
             "→ 生成 adapter_config.json + adapter_model.bin (~15MB)",
             "→ 推理时通过 PeftModel.from_pretrained() 动态加载",
         ], icon="💾")

add_card(slide, Inches(7.0), Inches(3.2), Inches(5.6), Inches(1.8),
         "方式二：合并到基础模型", [
             "merged = model.merge_and_unload()",
             "merged.save_pretrained(\"./merged_model\")",
             "",
             "→ 合并后结构与原始模型一致",
             "→ 推理无额外开销，适合生产部署",
         ], icon="🔗", title_color=COLOR_SECONDARY)

# 合并原理
add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.8), Inches(1.5),
          fill_color=RGBColor(0xF1, 0xF5, 0xF9))
add_text_box(slide, Inches(1.0), Inches(5.35), Inches(4), Inches(0.35),
             "合并原理", font_size=15, color=COLOR_PRIMARY, bold=True)
add_text_box(slide, Inches(1.0), Inches(5.75), Inches(11.4), Inches(0.9),
             "W_new = W₀ + (α / r) · B@A    其中 B@A 是矩阵乘法恢复原始维度\n"
             "合并后模型参数量与原始模型完全相同，前向传播无额外计算开销\n"
             "合并时需使用与训练一致的 dtype（如 bfloat16）以避免精度损失",
             font_size=14, color=COLOR_DARK)
add_page_number(slide, 17, TOTAL_SLIDES)


# ────────────────────────────────────
# S18: PART 6 章节页
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, fill_color=COLOR_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(3), Inches(0.6),
             "PART 06", font_size=16, color=COLOR_PRIMARY, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
             "评估方法与测试方案", font_size=42, color=COLOR_WHITE, bold=True)
add_shape(slide, Inches(1.5), Inches(3.9), Inches(3), Pt(3), fill_color=COLOR_ACCENT)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
             "自动化评估、人工评估、单元测试、集成测试、压力测试",
             font_size=18, color=COLOR_GRAY)
add_page_number(slide, 18, TOTAL_SLIDES)


# ────────────────────────────────────
# S19: 测试金字塔
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 6.1, "测试金字塔与评估指标")

# 金字塔 - 底层
pyramid_data = [
    (Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.55), "静态检查 —— 代码风格、类型标注、配置文件 Schema", RGBColor(0x93, 0xC5, 0xFD)),
    (Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.55), "单元测试 —— LoRA 配置、参数计算、矩阵维度、梯度检查", RGBColor(0x60, 0xA5, 0xFA)),
    (Inches(2.5), Inches(4.4), Inches(8.3), Inches(0.55), "集成测试 —— 训练管道、数据加载、模型保存/加载", RGBColor(0x3B, 0x82, 0xF6)),
    (Inches(3.5), Inches(5.0), Inches(6.3), Inches(0.55), "质量评估 —— ROUGE/BLEU/BERTScore, 人工评分", RGBColor(0x1D, 0x4E, 0xD8)),
    (Inches(4.5), Inches(5.6), Inches(4.3), Inches(0.55), "性能&压测 —— 延迟/吞吐/显存/并发", RGBColor(0x1E, 0x3A, 0x5A)),
    (Inches(5.0), Inches(6.2), Inches(3.3), Inches(0.55), "A/B 测试 —— 线上对比", RGBColor(0x0F, 0x1F, 0x3A)),
]

for left, top, width, height, label, color in pyramid_data:
    shape = add_shape(slide, left, top, width, height, fill_color=color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.05), width - Inches(0.4), height - Inches(0.1),
                 label, font_size=12, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 评估指标表
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(12), Inches(0.4),
             "核心指标: ROUGE-1/2/L | BLEU | BERTScore | Perplexity | 重复率 | 推理延迟 p50/p95/p99",
             font_size=13, color=COLOR_GRAY)
add_page_number(slide, 19, TOTAL_SLIDES)


# ────────────────────────────────────
# S20: PART 7 章节页
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, fill_color=COLOR_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(3), Inches(0.6),
             "PART 07", font_size=16, color=COLOR_PRIMARY, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
             "生产环境部署", font_size=42, color=COLOR_WHITE, bold=True)
add_shape(slide, Inches(1.5), Inches(3.9), Inches(3), Pt(3), fill_color=COLOR_ACCENT)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
             "vLLM / TGI / Docker / K8s 部署方案及生产监控",
             font_size=18, color=COLOR_GRAY)
add_page_number(slide, 20, TOTAL_SLIDES)


# ────────────────────────────────────
# S21: 部署方案对比
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 7.1, "部署方案对比")

deploy_data = [
    ("vLLM", "生产环境高并发", "低延迟", "✅ 原生多适配器", "🔥 推荐"),
    ("TGI", "HF 生态用户", "中等", "✅ 支持", ""),
    ("FastAPI+Trf", "快速原型", "中高", "❌ 需手动实现", ""),
    ("TensorRT-LLM", "极致性能", "极低延迟", "✅ 复杂配置", "🚀 极致优化"),
    ("llama.cpp", "CPU/边缘设备", "中等", "❌ 需 merge", ""),
]

y = Inches(3.5)
headers = ["方案", "适用场景", "延迟", "多适配器", "推荐"]
col_widths_d = [Inches(1.6), Inches(2.8), Inches(2.0), Inches(3.2), Inches(2.0)]
col_starts_d = [Inches(0.8)]
for w in col_widths_d[:-1]:
    col_starts_d.append(col_starts_d[-1] + w)

for j, (h, cs, cw) in enumerate(zip(headers, col_starts_d, col_widths_d)):
    add_text_box(slide, cs + Inches(0.1), y, cw - Inches(0.2), Inches(0.4),
                 h, font_size=13, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, cs, y, cw, Inches(0.4), fill_color=COLOR_PRIMARY)

for i, row_data in enumerate(deploy_data):
    yi = y + Inches(0.4) + i * Inches(0.42)
    bg = RGBColor(0xF8, 0xFA, 0xFC) if i % 2 == 0 else COLOR_WHITE
    for j, (cell, cs, cw) in enumerate(zip(row_data, col_starts_d, col_widths_d)):
        add_shape(slide, cs, yi, cw, Inches(0.42), fill_color=bg)
        cell_color = COLOR_ACCENT if j == 4 and "推荐" in cell else \
                     COLOR_ORANGE if j == 4 and "极致" in cell else \
                     COLOR_DARK
        cell_bold = j == 4
        add_text_box(slide, cs + Inches(0.1), yi + Inches(0.02), cw - Inches(0.2), Inches(0.38),
                     cell, font_size=12, color=cell_color, bold=cell_bold,
                     alignment=PP_ALIGN.CENTER)

# vLLM 代码
add_code_block(slide, Inches(0.8), Inches(6.0), Inches(11.8), Inches(1.2),
'''# vLLM 部署 (推荐)
vllm serve Qwen/Qwen2.5-7B-Instruct \
    --enable-lora \\
    --lora-modules adapter_a=./lora_a adapter_b=./lora_b \\
    --max-lora-rank=64 --port 8000''', font_size=11)
add_page_number(slide, 21, TOTAL_SLIDES)


# ────────────────────────────────────
# S22: PART 8 章节页
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, fill_color=COLOR_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(3), Inches(0.6),
             "PART 08", font_size=16, color=COLOR_PRIMARY, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
             "进阶话题与思考题", font_size=42, color=COLOR_WHITE, bold=True)
add_shape(slide, Inches(1.5), Inches(3.9), Inches(3), Pt(3), fill_color=COLOR_ACCENT)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
             "LoRA 变体演进、常见陷阱排查、启发性深度问题",
             font_size=18, color=COLOR_GRAY)
add_page_number(slide, 22, TOTAL_SLIDES)


# ────────────────────────────────────
# S23: LoRA 变体与陷阱
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_WHITE)
add_section_header(slide, 8.1, "LoRA 变体演进与常见陷阱")

# 变体
add_card(slide, Inches(0.8), Inches(3.2), Inches(5.8), Inches(3.5),
         "LoRA 变体家族", [
             "AdaLoRA (2023) — 自适应分配秩 (SVD分解)",
             "DoRA (2024) — 权重分解: magnitude + direction",
             "LoRA+ (2024) — A/B矩阵不同学习率",
             "PiSSA (2024) — 基于SVD的主成分初始化",
             "VeRA (2023) — 共享参数 + 可学习缩放",
             "Delta-LoRA (2024) — 参数增量放缩",
         ], icon="🔬")

# 常见陷阱
add_card(slide, Inches(7.0), Inches(3.2), Inches(5.6), Inches(3.5),
         "常见陷阱排查", [
             "❌ Loss 不降 → 检查 requires_grad, 增大LR",
             "❌ 生成重复内容 → LoRA过拟合, 降r/加dropout",
             "❌ 灾难性遗忘 → 数据混合: 10-20%通用数据",
             "❌ OOM → 减小 batch / 开启 gradient checkpoint",
             "❌ 适配器无效 → 用 model.print_model() 确认模块名",
             "❌ Merge 后效果下降 → 使用与训练一致 dtype",
         ], icon="⚠️", title_color=COLOR_RED)
add_page_number(slide, 23, TOTAL_SLIDES)


# ────────────────────────────────────
# S24: 结束页
# ────────────────────────────────────
slide = new_slide()
add_bg(slide, COLOR_DARK)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), fill_color=COLOR_PRIMARY)
add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), fill_color=COLOR_SECONDARY)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "感谢聆听", font_size=48, color=COLOR_WHITE, bold=True)
add_shape(slide, Inches(1), Inches(3.0), Inches(3), Pt(3), fill_color=COLOR_ACCENT)

add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             "LoRA 模型微调全流程讲解", font_size=24, color=COLOR_GRAY)

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.4),
             "✦  代码仓库: lora_demo.py （完整的可运行训练脚本）", font_size=16, color=COLOR_GRAY)
add_text_box(slide, Inches(1), Inches(4.95), Inches(11), Inches(0.4),
             "✦  配套文档: lora_finetuning_guide.md （12个启发性问题请见文档）", font_size=16, color=COLOR_GRAY)
add_text_box(slide, Inches(1), Inches(5.4), Inches(11), Inches(0.4),
             "✦  测试部署: test_deploy_guide.md （单元/集成/压测/部署全套方案）", font_size=16, color=COLOR_GRAY)

add_text_box(slide, Inches(1), Inches(6.5), Inches(5), Inches(0.4),
             "End", font_size=14, color=COLOR_GRAY)
add_page_number(slide, 24, TOTAL_SLIDES)


# ════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "LoRA_微调全流程讲解.pptx")
prs.save(output_path)
print(f"✅ PPTX 已生成: {output_path}")
print(f"   共 {len(prs.slides)} 页幻灯片")
